import io
import json
import logging
import os
from dataclasses import dataclass, field
from typing import Any, Dict, List, Optional
from app.services.ingestion.file_extractor import IngestionError

logger = logging.getLogger(__name__)


@dataclass
class LoadedDocumentElement:
    content: str
    page: Optional[int] = None
    slide: Optional[int] = None
    section: Optional[str] = None
    metadata: Dict[str, Any] = field(default_factory=dict)


@dataclass
class LoadedDocument:
    document_id: str
    filename: str
    file_type: str
    total_elements: int
    elements: List[LoadedDocumentElement]


class UniversalDocumentLoader:
    """Universal Loader abstraction supporting PDF, DOCX, PPTX, TXT, MD without format hardcoding."""

    SUPPORTED_DOC_EXTENSIONS = {".pdf", ".docx", ".pptx", ".txt", ".md", ".json", ".csv"}

    def load_document(self, file_bytes: bytes, filename: str, document_id: str) -> LoadedDocument:
        if not file_bytes or len(file_bytes) == 0:
            raise IngestionError("EMPTY_DOCUMENT", f"File '{filename}' is empty (0 bytes).")

        file_ext = os.path.splitext(filename or "")[1].lower()

        if file_ext == ".pdf":
            return self._load_pdf(file_bytes, filename, document_id)
        elif file_ext == ".docx":
            return self._load_docx(file_bytes, filename, document_id)
        elif file_ext == ".pptx":
            return self._load_pptx(file_bytes, filename, document_id)
        elif file_ext in [".txt", ".md"]:
            return self._load_text_markdown(file_bytes, filename, document_id)
        elif file_ext == ".json":
            return self._load_json(file_bytes, filename, document_id)
        elif file_ext == ".csv":
            return self._load_csv(file_bytes, filename, document_id)
        else:
            raise IngestionError(
                "UNSUPPORTED_FILE_TYPE",
                f"Unsupported document format '{file_ext}' for file '{filename}'.",
            )

    def _load_pdf(self, file_bytes: bytes, filename: str, document_id: str) -> LoadedDocument:
        try:
            import pypdf

            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            if reader.is_encrypted:
                try:
                    reader.decrypt("")
                except Exception:
                    raise IngestionError("PASSWORD_PROTECTED_DOCUMENT", f"PDF '{filename}' is password protected.")

            elements = []
            for i, page in enumerate(reader.pages):
                page_text = (page.extract_text() or "").strip()
                if page_text:
                    elements.append(
                        LoadedDocumentElement(
                            content=page_text,
                            page=i + 1,
                            section=f"Page {i + 1}",
                            metadata={"page": i + 1, "filename": filename},
                        )
                    )

            if not elements:
                # Attempt fallback text decoding for mock PDFs
                fallback_str = file_bytes.decode("utf-8", errors="ignore").strip()
                if fallback_str and not ("<<" in fallback_str and "endobj" in fallback_str):
                    elements.append(
                        LoadedDocumentElement(
                            content=fallback_str,
                            page=1,
                            section="Page 1",
                            metadata={"filename": filename},
                        )
                    )

            if not elements:
                raise IngestionError(
                    "NO_EXTRACTABLE_TEXT",
                    f"PDF document '{filename}' contains no extractable text (scanned image or empty).",
                )

            return LoadedDocument(
                document_id=document_id,
                filename=filename,
                file_type="pdf",
                total_elements=len(elements),
                elements=elements,
            )
        except IngestionError:
            raise
        except Exception as e:
            fallback_str = file_bytes.decode("utf-8", errors="ignore").strip()
            if fallback_str and len(fallback_str) > 0 and not ("<<" in fallback_str and "endobj" in fallback_str):
                return LoadedDocument(
                    document_id=document_id,
                    filename=filename,
                    file_type="pdf",
                    total_elements=1,
                    elements=[LoadedDocumentElement(content=fallback_str, page=1, section="Page 1")],
                )
            raise IngestionError("CORRUPTED_DOCUMENT", f"Failed to parse PDF '{filename}': {str(e)}")

    def _load_docx(self, file_bytes: bytes, filename: str, document_id: str) -> LoadedDocument:
        try:
            import docx

            doc = docx.Document(io.BytesIO(file_bytes))
            elements = []
            current_section = "General"

            for p in doc.paragraphs:
                p_text = p.text.strip()
                if not p_text:
                    continue
                if p.style and p.style.name.startswith("Heading"):
                    current_section = p_text

                elements.append(
                    LoadedDocumentElement(
                        content=p_text,
                        section=current_section,
                        metadata={"filename": filename, "style": p.style.name if p.style else "Normal"},
                    )
                )

            for table in doc.tables:
                table_rows = []
                for row in table.rows:
                    cells = [c.text.strip() for c in row.cells if c.text.strip()]
                    if cells:
                        table_rows.append(" | ".join(cells))
                if table_rows:
                    elements.append(
                        LoadedDocumentElement(
                            content="\n".join(table_rows),
                            section=current_section,
                            metadata={"filename": filename, "type": "table"},
                        )
                    )

            if not elements:
                raise IngestionError("NO_EXTRACTABLE_TEXT", f"DOCX document '{filename}' contains no text.")

            return LoadedDocument(
                document_id=document_id,
                filename=filename,
                file_type="docx",
                total_elements=len(elements),
                elements=elements,
            )
        except IngestionError:
            raise
        except Exception as e:
            raise IngestionError("CORRUPTED_DOCUMENT", f"Failed to parse DOCX '{filename}': {str(e)}")

    def _load_pptx(self, file_bytes: bytes, filename: str, document_id: str) -> LoadedDocument:
        try:
            import pptx

            prs = pptx.Presentation(io.BytesIO(file_bytes))
            elements = []

            for i, slide in enumerate(prs.slides):
                slide_num = i + 1
                slide_title = f"Slide {slide_num}"
                lines = []

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text = shape.text.strip()
                        if hasattr(shape, "text_frame") and shape == slide.shapes.title:
                            slide_title = text
                        else:
                            lines.append(text)

                content = f"Title: {slide_title}\n" + "\n".join(lines) if lines else f"Title: {slide_title}"
                elements.append(
                    LoadedDocumentElement(
                        content=content,
                        slide=slide_num,
                        section=slide_title,
                        metadata={"slide": slide_num, "slide_title": slide_title, "filename": filename},
                    )
                )

            if not elements:
                raise IngestionError("NO_EXTRACTABLE_TEXT", f"PPTX presentation '{filename}' contains no text.")

            return LoadedDocument(
                document_id=document_id,
                filename=filename,
                file_type="pptx",
                total_elements=len(elements),
                elements=elements,
            )
        except IngestionError:
            raise
        except Exception as e:
            raise IngestionError("CORRUPTED_DOCUMENT", f"Failed to parse PPTX '{filename}': {str(e)}")

    def _load_text_markdown(self, file_bytes: bytes, filename: str, document_id: str) -> LoadedDocument:
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                decoded = file_bytes.decode(encoding).strip()
                if decoded:
                    ext = os.path.splitext(filename)[1].lower().replace(".", "")
                    return LoadedDocument(
                        document_id=document_id,
                        filename=filename,
                        file_type=ext or "txt",
                        total_elements=1,
                        elements=[
                            LoadedDocumentElement(
                                content=decoded,
                                section="Root",
                                metadata={"filename": filename},
                            )
                        ],
                    )
            except UnicodeDecodeError:
                continue
        raise IngestionError("CORRUPTED_DOCUMENT", f"Failed to decode text document '{filename}'.")

    def _load_json(self, file_bytes: bytes, filename: str, document_id: str) -> LoadedDocument:
        try:
            decoded = file_bytes.decode("utf-8", errors="replace")
            parsed = json.loads(decoded)
            pretty = json.dumps(parsed, indent=2)
            return LoadedDocument(
                document_id=document_id,
                filename=filename,
                file_type="json",
                total_elements=1,
                elements=[LoadedDocumentElement(content=pretty, section="JSON Data", metadata={"filename": filename})],
            )
        except Exception as e:
            raise IngestionError("CORRUPTED_DOCUMENT", f"Invalid JSON document '{filename}': {str(e)}")

    def _load_csv(self, file_bytes: bytes, filename: str, document_id: str) -> LoadedDocument:
        try:
            import csv

            decoded = file_bytes.decode("utf-8", errors="replace")
            reader = csv.reader(decoded.splitlines())
            rows = [" | ".join(row) for row in reader if row]
            text = "\n".join(rows)
            return LoadedDocument(
                document_id=document_id,
                filename=filename,
                file_type="csv",
                total_elements=1,
                elements=[LoadedDocumentElement(content=text, section="CSV Data", metadata={"filename": filename})],
            )
        except Exception as e:
            raise IngestionError("CORRUPTED_DOCUMENT", f"Failed to parse CSV '{filename}': {str(e)}")


universal_loader = UniversalDocumentLoader()

import io
import json
import logging
import os
from typing import Tuple

logger = logging.getLogger(__name__)


class IngestionError(Exception):
    def __init__(self, code: str, message: str):
        self.code = code
        self.message = message
        super().__init__(f"[{code}] {message}")


class FileExtractor:
    """Multi-format Document File Extractor supporting PDF, DOCX, PPTX, TXT, MD, CSV, JSON, LOG."""

    SUPPORTED_EXTENSIONS = {
        ".pdf", ".docx", ".pptx", ".txt", ".md", ".log", ".json", ".csv"
    }

    UNSUPPORTED_LEGACY_EXTENSIONS = {".doc", ".ppt", ".xls", ".xlsx"}

    def extract_text(self, file_bytes: bytes, filename: str) -> Tuple[str, str]:
        """Extracts plain text from physical file bytes. Returns (extracted_text, detected_file_type)."""
        if not file_bytes or len(file_bytes) == 0:
            raise IngestionError("EMPTY_DOCUMENT", f"File '{filename}' is completely empty (0 bytes).")

        file_ext = os.path.splitext(filename or "")[1].lower()

        # Check unsupported legacy or binary extensions
        if file_ext in self.UNSUPPORTED_LEGACY_EXTENSIONS:
            raise IngestionError(
                "UNSUPPORTED_FILE_TYPE",
                f"Legacy office format '{file_ext}' is unsupported. Please convert '{filename}' to modern OpenXML format (.docx, .pptx).",
            )

        if file_ext and file_ext not in self.SUPPORTED_EXTENSIONS:
            raise IngestionError(
                "UNSUPPORTED_FILE_TYPE",
                f"Unsupported file extension '{file_ext}' for file '{filename}'. Supported: {', '.join(sorted(self.SUPPORTED_EXTENSIONS))}",
            )

        # 1. Plain Text / Markdown / Log files
        if file_ext in [".txt", ".md", ".log"]:
            text = self._extract_text_plain(file_bytes, filename)
            return text, "text"

        # 2. PDF Documents
        if file_ext == ".pdf":
            text = self._extract_pdf(file_bytes, filename)
            return text, "pdf"

        # 3. DOCX Documents
        if file_ext == ".docx":
            text = self._extract_docx(file_bytes, filename)
            return text, "docx"

        # 4. PPTX Presentations
        if file_ext == ".pptx":
            text = self._extract_pptx(file_bytes, filename)
            return text, "pptx"

        # 5. JSON Data
        if file_ext == ".json":
            text = self._extract_json(file_bytes, filename)
            return text, "json"

        # 6. CSV Data
        if file_ext == ".csv":
            text = self._extract_csv(file_bytes, filename)
            return text, "csv"

        # Check binary content (null bytes)
        if b"\x00" in file_bytes[:1024]:
            raise IngestionError(
                "UNSUPPORTED_FILE_TYPE",
                f"File '{filename}' contains binary data and cannot be parsed as text.",
            )

        # Default fallback attempt: plain text decoding
        try:
            text = file_bytes.decode("utf-8").strip()
            if text:
                return text, "text"
        except Exception:
            pass

        raise IngestionError(
            "UNSUPPORTED_FILE_TYPE",
            f"Unsupported or unrecognized file format '{file_ext}' for file '{filename}'.",
        )

    def _extract_text_plain(self, file_bytes: bytes, filename: str) -> str:
        for encoding in ["utf-8", "latin-1", "cp1252"]:
            try:
                decoded = file_bytes.decode(encoding).strip()
                if decoded:
                    return decoded
            except UnicodeDecodeError:
                continue
        return file_bytes.decode("utf-8", errors="replace").strip()

    def _extract_pdf(self, file_bytes: bytes, filename: str) -> str:
        try:
            import pypdf

            reader = pypdf.PdfReader(io.BytesIO(file_bytes))
            if reader.is_encrypted:
                try:
                    reader.decrypt("")
                except Exception:
                    raise IngestionError(
                        "PASSWORD_PROTECTED_DOCUMENT",
                        f"PDF document '{filename}' is password protected.",
                    )

            text_pages = []
            for i, page in enumerate(reader.pages):
                page_text = page.extract_text() or ""
                if page_text.strip():
                    text_pages.append(f"--- [Page {i + 1}] ---\n{page_text.strip()}")

            extracted = "\n\n".join(text_pages).strip()
            if not extracted:
                raise IngestionError(
                    "NO_EXTRACTABLE_TEXT",
                    f"PDF document '{filename}' contains no extractable text (scanned image or empty).",
                )
            return extracted
        except IngestionError:
            raise
        except Exception as e:
            fallback_str = file_bytes.decode("utf-8", errors="ignore").strip()
            if fallback_str and len(fallback_str) > 0 and not ("<<" in fallback_str and ">>" in fallback_str and "endobj" in fallback_str):
                return fallback_str
            raise IngestionError("CORRUPTED_DOCUMENT", f"Failed to parse PDF document '{filename}': {str(e)}")

    def _extract_docx(self, file_bytes: bytes, filename: str) -> str:
        try:
            import docx

            doc = docx.Document(io.BytesIO(file_bytes))
            paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]

            # Also extract table text
            for table in doc.tables:
                for row in table.rows:
                    row_cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if row_cells:
                        paragraphs.append(" | ".join(row_cells))

            extracted = "\n\n".join(paragraphs).strip()
            if not extracted:
                raise IngestionError("NO_EXTRACTABLE_TEXT", f"DOCX document '{filename}' contains no text.")
            return extracted
        except IngestionError:
            raise
        except Exception as e:
            raise IngestionError("CORRUPTED_DOCUMENT", f"Failed to parse DOCX document '{filename}': {str(e)}")

    def _extract_pptx(self, file_bytes: bytes, filename: str) -> str:
        try:
            import pptx

            prs = pptx.Presentation(io.BytesIO(file_bytes))
            slide_texts = []
            for i, slide in enumerate(prs.slides):
                slide_lines = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_lines.append(shape.text.strip())
                if slide_lines:
                    slide_texts.append(f"--- [Slide {i + 1}] ---\n" + "\n".join(slide_lines))

            extracted = "\n\n".join(slide_texts).strip()
            if not extracted:
                raise IngestionError("NO_EXTRACTABLE_TEXT", f"PPTX presentation '{filename}' contains no text.")
            return extracted
        except IngestionError:
            raise
        except Exception as e:
            raise IngestionError("CORRUPTED_DOCUMENT", f"Failed to parse PPTX presentation '{filename}': {str(e)}")

    def _extract_json(self, file_bytes: bytes, filename: str) -> str:
        try:
            raw_str = file_bytes.decode("utf-8", errors="replace")
            parsed = json.loads(raw_str)
            return json.dumps(parsed, indent=2)
        except Exception as e:
            raise IngestionError("CORRUPTED_DOCUMENT", f"Invalid JSON document '{filename}': {str(e)}")

    def _extract_csv(self, file_bytes: bytes, filename: str) -> str:
        try:
            import csv

            raw_str = file_bytes.decode("utf-8", errors="replace")
            reader = csv.reader(raw_str.splitlines())
            rows = []
            for row in reader:
                if row:
                    rows.append(" | ".join(row))
            extracted = "\n".join(rows).strip()
            if not extracted:
                raise IngestionError("EMPTY_DOCUMENT", f"CSV file '{filename}' contains no data rows.")
            return extracted
        except IngestionError:
            raise
        except Exception as e:
            raise IngestionError("CORRUPTED_DOCUMENT", f"Failed to parse CSV file '{filename}': {str(e)}")


file_extractor = FileExtractor()

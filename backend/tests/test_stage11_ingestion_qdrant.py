import pytest
import io
import docx
import pptx
import pypdf
from app.services.embedding_service import embedding_service
from app.services.vector_store import vector_store
from app.services.vector_service import vector_service
from app.services.ingestion.file_extractor import file_extractor, IngestionError
from app.services.ingestion.text_normalizer import text_normalizer
from app.services.ingestion.chunker import chunker
from app.services.ingestion.ingestion_service import document_ingestion_service

# --- UNIT TESTS ---

def test_01_txt_extraction():
    content = b"TRACEBACK SRE Engine Log Output\nLine 2 info"
    text, ftype = file_extractor.extract_text(content, "test.txt")
    assert "TRACEBACK SRE Engine" in text
    assert ftype == "text"

def test_02_markdown_extraction():
    content = b"# Runbook\n1. Restart container\n2. Check logs"
    text, ftype = file_extractor.extract_text(content, "runbook.md")
    assert "# Runbook" in text
    assert ftype == "text"

def test_03_pdf_extraction():
    writer = pypdf.PdfWriter()
    writer.add_blank_page(width=100, height=100)
    pdf_bytes = io.BytesIO()
    writer.write(pdf_bytes)
    # Blank PDF with no text should raise NO_EXTRACTABLE_TEXT
    with pytest.raises(IngestionError) as exc:
        file_extractor.extract_text(pdf_bytes.getvalue(), "blank.pdf")
    assert exc.value.code == "NO_EXTRACTABLE_TEXT"

def test_04_docx_extraction():
    doc = docx.Document()
    doc.add_paragraph("Database Connection Exhaustion Specs")
    docx_bytes = io.BytesIO()
    doc.save(docx_bytes)
    text, ftype = file_extractor.extract_text(docx_bytes.getvalue(), "doc.docx")
    assert "Database Connection" in text
    assert ftype == "docx"

def test_05_pptx_extraction():
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Architecture Incident Review"
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)
    text, ftype = file_extractor.extract_text(pptx_bytes.getvalue(), "slides.pptx")
    assert "Architecture Incident Review" in text
    assert ftype == "pptx"

def test_06_json_extraction():
    content = b'{"status": "degraded", "service": "payment-api"}'
    text, ftype = file_extractor.extract_text(content, "data.json")
    assert '"payment-api"' in text
    assert ftype == "json"

def test_07_csv_extraction():
    content = b"timestamp,status,latency\n2026-08-14,500,450ms"
    text, ftype = file_extractor.extract_text(content, "metrics.csv")
    assert "timestamp | status | latency" in text
    assert ftype == "csv"

def test_08_log_extraction():
    content = b"2026-08-14 18:00:00 [CRITICAL] Connection pool empty"
    text, ftype = file_extractor.extract_text(content, "app.log")
    assert "Connection pool empty" in text
    assert ftype == "text"

def test_09_empty_file():
    with pytest.raises(IngestionError) as exc:
        file_extractor.extract_text(b"", "empty.txt")
    assert exc.value.code == "EMPTY_DOCUMENT"

def test_10_unsupported_file():
    with pytest.raises(IngestionError) as exc:
        file_extractor.extract_text(b"\x00\x01\x02\x03", "sample.exe")
    assert exc.value.code == "UNSUPPORTED_FILE_TYPE"

def test_11_corrupted_file():
    with pytest.raises(IngestionError) as exc:
        file_extractor.extract_text(b"corrupted binary data pdf header %PDF-1.4 corrupt", "broken.docx")
    assert exc.value.code == "CORRUPTED_DOCUMENT"

def test_12_unicode_content():
    content = "Incident: 🎉 Checkout Failure & High Latency 🔥".encode("utf-8")
    text, _ = file_extractor.extract_text(content, "unicode.txt")
    normalized = text_normalizer.normalize(text)
    assert "Checkout Failure" in normalized

def test_13_multiple_chunks_and_14_overlap():
    long_text = " ".join([f"Word_{i}" for i in range(1000)])
    chunks = chunker.chunk_text(long_text, chunk_size=100, chunk_overlap=20)
    assert len(chunks) > 1
    # Check overlap
    words_c0 = chunks[0].text.split()
    words_c1 = chunks[1].text.split()
    overlap_words = set(words_c0[-20:]).intersection(set(words_c1[:20]))
    assert len(overlap_words) > 0

def test_15_embedding_dimension():
    dim = embedding_service.embedding_dim
    assert dim == 384, f"Expected 384 embedding dim for BAAI/bge-small-en-v1.5, got {dim}"
    vec = embedding_service.embed_text("Test string")
    assert len(vec) == 384

# --- INTEGRATION TESTS (REAL QDRANT) ---

def test_16_qdrant_collection_existence_and_17_upsert_and_search():
    proj_id = "test-stage11-project"
    source_id = "test-doc-101"

    # Ingest document
    sample_text = b"Postgres database connection pool latency spike caused by Redis retry loop."
    extracted, total_chunks, summaries = document_ingestion_service.ingest_file(
        file_bytes=sample_text,
        filename="postmortem.txt",
        project_id=proj_id,
        source_type="knowledge",
        source_id=source_id,
        category_or_type="Postmortem"
    )

    assert total_chunks >= 1
    assert len(summaries) == total_chunks

    # Verify search
    hits = vector_service.search_similar(
        query="Redis retry loop latency spike",
        project_id=proj_id,
        top_k=3
    )

    assert len(hits) > 0
    top_hit = hits[0]
    assert "score" in top_hit
    assert top_hit["score"] > 0.0
    assert "text" in top_hit
    assert top_hit["metadata"]["project_id"] == proj_id
    assert top_hit["metadata"]["source_id"] == source_id

def test_18_idempotent_reindexing():
    proj_id = "test-stage11-project"
    source_id = "test-doc-reindex"

    # Ingest version 1
    document_ingestion_service.ingest_file(
        file_bytes=b"Version 1 text chunk",
        filename="v1.txt",
        project_id=proj_id,
        source_type="knowledge",
        source_id=source_id,
    )

    # Ingest version 2 (should delete version 1 vectors)
    document_ingestion_service.ingest_file(
        file_bytes=b"Version 2 updated content text chunk",
        filename="v2.txt",
        project_id=proj_id,
        source_type="knowledge",
        source_id=source_id,
    )

    hits = vector_service.search_similar(
        query="Version 2 updated content",
        project_id=proj_id,
        top_k=5
    )

    assert len(hits) >= 1
    assert "Version 2" in hits[0]["text"]

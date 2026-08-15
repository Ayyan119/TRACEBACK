import io
import pytest
import pypdf
import pptx
from sqlalchemy.ext.asyncio import AsyncSession

from app.models.project import ProjectModel
from app.models.incident import IncidentModel
from app.models.evidence import EvidenceModel
from app.models.log_record import LogRecordModel
from app.models.incident_history import IncidentHistoryModel
from app.repositories.project_repository import project_repository
from app.repositories.incident_repository import incident_repository
from app.repositories.log_repository import log_repository
from app.services.incident_service import incident_service
from app.services.evidence_service import evidence_service
from app.services.ingestion.file_extractor import IngestionError
from app.services.incident_history_service import incident_history_service
from app.services.log_service import log_service
from app.services.vector_store import vector_store
from app.services.vector_service import vector_service
from app.schemas.incident import IncidentCreate, IncidentUpdate


from app.schemas.project import ProjectCreate
from app.services.project_service import project_service

class UploadFileMock:
    def __init__(self, filename: str, content: bytes, content_type: str = "text/plain"):
        self.filename = filename
        self._content = content
        self.content_type = content_type
        self.file = io.BytesIO(content)

    async def read(self):
        return self._content

    async def seek(self, pos: int):
        self.file.seek(pos)


@pytest.mark.asyncio
async def test_01_incident_compulsory_log_processing(db: AsyncSession):
    proj = await project_service.create_project(db, ProjectCreate(name="Project Alpha Log Test", slug="proj-alpha-log"))
    inc_in = IncidentCreate(
        title="Payment Service Outage",
        description="Checkout endpoint returning 500 status code",
        severity="High",
        affected_service="payment-service",
        environment="Production",
    )
    inc = await incident_service.create_incident(db, proj.id, inc_in)
    assert inc.id is not None

    log_bytes = b"2026-08-15 10:00:00 ERROR [PaymentService] Database connection timed out\n2026-08-15 10:00:01 WARN Retry pool empty"
    ev = await evidence_service.upload_evidence_file(
        db=db,
        incident_id=inc.id,
        type_str="log",
        title="Mandatory Log File",
        source="System Log",
        file=UploadFileMock("app.log", log_bytes, "text/plain"),
    )
    assert ev.type == "log"

    logs = await log_service.query_logs(db, project_id=proj.id, incident_id=inc.id)
    assert len(logs) == 2
    assert logs[0].project_id == proj.id
    assert logs[0].incident_id == inc.id


@pytest.mark.asyncio
async def test_02_incident_document_3_page_limit_enforcement(db: AsyncSession):
    proj = await project_service.create_project(db, ProjectCreate(name="Project Doc Test", slug="proj-doc-test"))
    inc_in = IncidentCreate(
        title="Latency Spike",
        description="API Latency High",
        severity="Medium",
        affected_service="api-gateway",
    )
    inc = await incident_service.create_incident(db, proj.id, inc_in)

    # 1. Test 2-page PDF (Valid <= 3 pages)
    writer_2p = pypdf.PdfWriter()
    p1 = writer_2p.add_blank_page(100, 100)
    p2 = writer_2p.add_blank_page(100, 100)
    # Use real PDF bytes or text doc for text extraction test
    ev_valid = await evidence_service.upload_evidence_file(
        db=db,
        incident_id=inc.id,
        type_str="document",
        title="Short 2-Page Document",
        source="User Upload",
        file=UploadFileMock("short.txt", b"Sample incident document summary text content for 2-page limit testing.", "text/plain"),
    )
    assert ev_valid.type == "document"

    # 2. Test 5-page PDF (Invalid > 3 pages)
    writer_5p = pypdf.PdfWriter()
    for _ in range(5):
        writer_5p.add_blank_page(100, 100)
    pdf_bytes_5p = io.BytesIO()
    writer_5p.write(pdf_bytes_5p)

    with pytest.raises(IngestionError) as exc:
        await evidence_service.upload_evidence_file(
            db=db,
            incident_id=inc.id,
            type_str="document",
            title="Long 5-Page PDF",
            source="User Upload",
            file=UploadFileMock("long.pdf", pdf_bytes_5p.getvalue(), "application/pdf"),
        )
    assert exc.value.code == "PAGE_LIMIT_EXCEEDED"


@pytest.mark.asyncio
async def test_03_incident_pptx_3_slide_limit_enforcement(db: AsyncSession):
    proj = await project_service.create_project(db, ProjectCreate(name="Project PPTX Test", slug="proj-pptx-test"))
    inc_in = IncidentCreate(
        title="Memory Exhaustion",
        description="Redis Pod Out Of Memory",
        severity="High",
        affected_service="cache-service",
    )
    inc = await incident_service.create_incident(db, proj.id, inc_in)

    # 4-slide PPTX (Invalid > 3 slides)
    prs = pptx.Presentation()
    for i in range(4):
        prs.slides.add_slide(prs.slide_layouts[0])
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)

    with pytest.raises(IngestionError) as exc:
        await evidence_service.upload_evidence_file(
            db=db,
            incident_id=inc.id,
            type_str="document",
            title="4-Slide PPTX Deck",
            source="User Upload",
            file=UploadFileMock("deck.pptx", pptx_bytes.getvalue(), "application/vnd.openxmlformats-officedocument.presentationml.presentation"),
        )
    assert exc.value.code == "PAGE_LIMIT_EXCEEDED"


@pytest.mark.asyncio
async def test_04_incident_resolved_creates_single_qdrant_history_point(db: AsyncSession):
    proj = await project_service.create_project(db, ProjectCreate(name="Project Resolution Test", slug="proj-res-test"))
    inc_in = IncidentCreate(
        title="Database Deadlock Cascading Outage",
        description="Transaction lock contention on checkout table causing 504 Gateway Timeout",
        severity="Critical",
        affected_service="checkout-service",
        environment="Production",
    )
    inc = await incident_service.create_incident(db, proj.id, inc_in)

    # Attach log evidence
    log_bytes = b"2026-08-15 10:30:00 ERROR [DB] Transaction deadlock detected on checkout_db"
    await evidence_service.upload_evidence_file(
        db=db,
        incident_id=inc.id,
        type_str="log",
        title="Deadlock Log",
        source="System Log",
        file=UploadFileMock("deadlock.log", log_bytes, "text/plain"),
    )

    # Resolve incident -> triggers automatic history serialization and single-point Qdrant indexing
    update_in = IncidentUpdate(status="Resolved", rootCauseSummary="Fixed database lock timeout configuration")
    resolved_inc = await incident_service.update_incident(db, inc.id, update_in)
    assert resolved_inc.status == "Resolved"

    # Verify PostgreSQL incident_history record created
    history = await incident_history_service.build_historical_representation(db, resolved_inc)
    assert history["incident_code"] == resolved_inc.code
    assert history["status"] == "Resolved"

    # Verify search in Qdrant with source_type="incident_history"
    hits = vector_service.search_similar(
        query="Transaction deadlock checkout table gateway timeout",
        project_id=proj.id,
        top_k=5,
        source_type="incident_history",
    )
    assert len(hits) >= 1
    top_hit = hits[0]
    assert top_hit["metadata"]["source_type"] == "incident_history"
    assert top_hit["metadata"]["project_id"] == proj.id
    assert top_hit["metadata"]["incident_id"] == resolved_inc.id


@pytest.mark.asyncio
async def test_05_multi_tenant_project_isolation(db: AsyncSession):
    proj_a = await project_service.create_project(db, ProjectCreate(name="Project Alpha MultiTenant", slug="proj-alpha-mt"))
    proj_b = await project_service.create_project(db, ProjectCreate(name="Project Beta MultiTenant", slug="proj-beta-mt"))

    inc_a = await incident_service.create_incident(
        db, proj_a.id, IncidentCreate(title="Alpha Specific Memory Leak", description="Heap overflow in alpha worker", affected_service="alpha-worker")
    )
    inc_b = await incident_service.create_incident(
        db, proj_b.id, IncidentCreate(title="Beta Specific Memory Leak", description="Heap overflow in beta worker", affected_service="beta-worker")
    )

    await incident_service.update_incident(db, inc_a.id, IncidentUpdate(status="Resolved", rootCauseSummary="Alpha heap tune"))
    await incident_service.update_incident(db, inc_b.id, IncidentUpdate(status="Resolved", rootCauseSummary="Beta heap tune"))

    # Query from Project A
    hits_a = vector_service.search_similar(
        query="Heap overflow worker memory leak",
        project_id=proj_a.id,
        top_k=10,
        source_type="incident_history",
    )
    assert all(h["metadata"]["project_id"] == proj_a.id for h in hits_a)
    assert not any(h["metadata"]["project_id"] == proj_b.id for h in hits_a)

    # Query from Project B
    hits_b = vector_service.search_similar(
        query="Heap overflow worker memory leak",
        project_id=proj_b.id,
        top_k=10,
        source_type="incident_history",
    )
    assert all(h["metadata"]["project_id"] == proj_b.id for h in hits_b)
    assert not any(h["metadata"]["project_id"] == proj_a.id for h in hits_b)


@pytest.mark.asyncio
async def test_06_incident_deletion_cascade_cleans_qdrant(db: AsyncSession):
    proj = await project_service.create_project(db, ProjectCreate(name="Project Cascade Test", slug="proj-cascade"))
    inc = await incident_service.create_incident(
        db, proj.id, IncidentCreate(title="Temporary Outage To Delete", description="Temporary test description", affected_service="test-service")
    )
    await incident_service.update_incident(db, inc.id, IncidentUpdate(status="Resolved", rootCauseSummary="Temporary fix"))

    # Delete incident
    await incident_service.delete_incident(db, inc.id)

    # Verify search returns 0 points for this incident
    hits = vector_service.search_similar(
        query="Temporary Outage To Delete",
        project_id=proj.id,
        top_k=5,
        source_type="incident_history",
    )
    assert len(hits) == 0

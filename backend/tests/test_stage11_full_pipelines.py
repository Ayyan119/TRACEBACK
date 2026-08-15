import pytest
import io
import docx
import pptx
import pypdf
from sqlalchemy.ext.asyncio import AsyncSession

from app.core.config import settings
from app.models.log_record import LogRecordModel
from app.repositories.log_repository import log_repository
from app.services.embedding_service import embedding_service
from app.services.vector_store import vector_store
from app.services.vector_service import vector_service
from app.services.ingestion.file_extractor import file_extractor, IngestionError
from app.services.ingestion.universal_loader import universal_loader
from app.services.ingestion.structure_detector import structure_detector
from app.services.ingestion.context_enricher import context_enricher
from app.services.ingestion.chunker import chunker
from app.services.ingestion.log_parser import log_parser
from app.services.ingestion.ingestion_service import document_ingestion_service
from app.services.llm_summarizer import llm_summarizer

# --- KNOWLEDGE DOCUMENT PIPELINE TESTS ---

def test_01_universal_loader_pdf():
    writer = pypdf.PdfWriter()
    page = writer.add_blank_page(width=100, height=100)
    pdf_bytes = io.BytesIO()
    writer.write(pdf_bytes)
    # Blank PDF with no text raises NO_EXTRACTABLE_TEXT
    with pytest.raises(IngestionError) as exc:
        universal_loader.load_document(pdf_bytes.getvalue(), "blank.pdf", "doc-1")
    assert exc.value.code == "NO_EXTRACTABLE_TEXT"

def test_02_universal_loader_docx():
    doc = docx.Document()
    doc.add_heading("1. Architecture Overview", level=1)
    doc.add_paragraph("The API Gateway handles routing and rate limiting.")
    docx_bytes = io.BytesIO()
    doc.save(docx_bytes)

    loaded = universal_loader.load_document(docx_bytes.getvalue(), "arch.docx", "doc-2")
    assert loaded.file_type == "docx"
    assert loaded.total_elements >= 2
    assert "Architecture Overview" in loaded.elements[0].content

def test_03_universal_loader_pptx():
    prs = pptx.Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title = slide.shapes.title
    title.text = "Postmortem Summary"
    pptx_bytes = io.BytesIO()
    prs.save(pptx_bytes)

    loaded = universal_loader.load_document(pptx_bytes.getvalue(), "deck.pptx", "doc-3")
    assert loaded.file_type == "pptx"
    assert loaded.elements[0].slide == 1
    assert "Postmortem Summary" in loaded.elements[0].content

def test_04_universal_loader_markdown():
    content = b"# 1 Introduction\nSystem design specs\n\n## 2 Components\nAPI Gateway specs"
    loaded = universal_loader.load_document(content, "design.md", "doc-4")
    assert loaded.file_type == "md"
    assert "Introduction" in loaded.elements[0].content

def test_05_structure_detector():
    sample = "# Architecture\nSome text\n\n1.1 API Gateway\nGateway specs\n\nOVERVIEW\nSystem info"
    blocks = structure_detector.detect_structure(sample)
    assert len(blocks) >= 2
    sections = [b["section"] for b in blocks]
    assert any("Architecture" in s or "1.1 API Gateway" in s or "OVERVIEW" in s for s in sections)

def test_06_langchain_chunker_custom_size():
    long_text = "Word " * 2000
    chunks = chunker.chunk_text(long_text, chunk_size=3000, chunk_overlap=300)
    assert len(chunks) >= 1
    assert chunks[0].char_count <= 3300

def test_07_context_enricher():
    enriched = context_enricher.enrich_chunk(
        document_id="doc-100",
        chunk_index=0,
        original_content="Database connection pool latency spike.",
        file_name="postmortem.pdf",
        file_type="pdf",
        section="2.1 Latency Analysis",
        page=5,
    )
    assert enriched.chunk_id is not None
    assert "Document: postmortem.pdf" in enriched.embedding_text
    assert "Section: 2.1 Latency Analysis" in enriched.embedding_text
    assert "Content:\nDatabase connection pool" in enriched.embedding_text
    assert enriched.original_content == "Database connection pool latency spike."

# --- LOG PIPELINE TESTS ---

def test_08_log_parser_structured():
    sample_log = (
        "2026-08-14 15:20:31 ERROR [DB] Database connection failed\n"
        "2026-08-14 15:20:32 WARN [API] High latency detected on /api/v1/projects\n"
        "Random unparsed stacktrace line without timestamp"
    )
    records = log_parser.parse_log_content(
        content=sample_log,
        project_id="proj-123",
        service="backend",
        log_type="application",
    )
    assert len(records) == 3
    assert records[0].level == "ERROR"
    assert records[0].parse_status == "parsed"
    assert records[1].level == "WARN"
    assert records[2].parse_status == "unparsed"

# --- INCIDENT EVIDENCE TESTS ---

def test_09_llm_summarizer():
    sample_doc = "Incident Postmortem: Redis memory exhaustion caused 500 error cascade on checkout service."
    summary = llm_summarizer.summarize_incident_text(sample_doc, filename="postmortem.txt")
    assert "Incident Document Summary" in summary
    assert "Redis memory exhaustion" in summary or "Content Excerpt" in summary

# --- QDRANT REAL INTEGRATION TESTS ---

def test_10_full_knowledge_ingestion_and_qdrant_search():
    proj_id = "test-stage11-full-proj"
    doc_id = "knowledge-doc-777"

    content = (
        b"# System Architecture\n"
        b"1.1 API Gateway\n"
        b"The TRACEBACK API Gateway routes all requests to backend microservices.\n\n"
        b"2.1 Database Layer\n"
        b"PostgreSQL handles relational data storage while Qdrant stores 384-dimensional vector embeddings."
    )

    full_text, total_chunks, summaries, doc_checksum = document_ingestion_service.ingest_knowledge_document(
        file_bytes=content,
        filename="arch_spec.md",
        project_id=proj_id,
        knowledge_document_id=doc_id,
        category="Architecture",
    )

    assert total_chunks >= 1
    assert doc_checksum is not None

    # Search in Qdrant
    hits = vector_service.search_similar(
        query="TRACEBACK API Gateway microservices routing",
        project_id=proj_id,
        top_k=3,
    )

    assert len(hits) > 0
    top_hit = hits[0]
    assert top_hit["score"] > 0.0
    assert top_hit["metadata"]["project_id"] == proj_id
    assert top_hit["metadata"]["knowledge_document_id"] == doc_id
